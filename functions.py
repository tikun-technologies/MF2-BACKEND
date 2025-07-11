import datetime
import requests
from selenium import webdriver
import time
import json
from selenium.webdriver.common.by import By
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pandas as pd
import tempfile
import shutil
from pptx.enum.shapes import MSO_SHAPE
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from PIL import Image
from pptx.enum.text import MSO_AUTO_SIZE

from selenium.webdriver.chrome.options import Options

from selenium.webdriver.common.action_chains import ActionChains


import json
import uuid
import pandas as pd
from io import BytesIO
from functools import wraps
from DB.db import STUDIES_collection
from flask_jwt_extended import  jwt_required
from azure.storage.blob import BlobServiceClient

AZURE_STORAGE_CONNECTION_STRING = "DefaultEndpointsProtocol=https;AccountName=printxd;AccountKey=CaL/3SmhK8iKVM02i/cIN1VgE3058lyxRnCxeRd2J1k/9Ay6I67GC2CMnW//lJhNl+71WwxYXHnC+AStkbW1Jg==;EndpointSuffix=core.windows.net"
CONTAINER_NAME = "mf2"

blob_service_client = BlobServiceClient.from_connection_string(AZURE_STORAGE_CONNECTION_STRING)
def upload_to_azure(file):
    blob_name = str(uuid.uuid4()) + "_" + file.filename  # Unique filename
    blob_client = blob_service_client.get_blob_client(container=CONTAINER_NAME, blob=blob_name)
    
    blob_client.upload_blob(file, overwrite=True)
    return blob_client.url  # Return file URL

def upload_to_azure_with_filename(file,filename):
    blob_name = str(uuid.uuid4()) + "_" + filename  # Unique filename
    blob_client = blob_service_client.get_blob_client(container=CONTAINER_NAME, blob=blob_name)
    
    blob_client.upload_blob(file, overwrite=True)
    return blob_client.url  # Return file URL
def protected(f):
    @wraps(f)
    @jwt_required()
    def wrapper(*args, **kwargs):
        return f(*args, **kwargs)
    return wrapper



def extract_summarizer_data(file_bytes):
    # Load the Excel file
    xls = pd.ExcelFile(BytesIO(file_bytes))
    
    # Identify summarizer sheets
    summarizer_sheets = [sheet for sheet in xls.sheet_names if "Summarizer" in sheet]
    
    # Function to process a single sheet
    def process_sheet(df):
        formatted_data = {}
        current_heading = None
        current_content = []
        headings = list(df[0].dropna())
        
        if headings:
            first_heading = headings[0].strip()
            formatted_data[first_heading] = ""
            first_heading_started = False
        
        for _, row in df.iterrows():
            cell_value = str(row[0]).strip()
            
            if not cell_value or cell_value.lower() == "nan":
                continue  # Skip empty rows
            
            # Detect headings (bold-like) and content
            if cell_value.isupper() or cell_value.endswith(":"):
                if current_heading and current_content:
                    formatted_data[current_heading] = "\n".join(current_content)
                current_heading = cell_value.replace(":", "").strip()
                current_content = []
                first_heading_started = True
            else:
                if not first_heading_started and first_heading:
                    formatted_data[first_heading] += f"{cell_value} "
                else:
                    current_content.append(cell_value)
        
        # Add last section
        if current_heading and current_content:
            formatted_data[current_heading] = "\n".join(current_content)
        
        return formatted_data
    
    # Extract and process all summarizer sheets
    summarizer_data = {sheet: process_sheet(xls.parse(sheet, header=None)) for sheet in summarizer_sheets}
    return summarizer_data


def get_file_data_for_study(file_bytes):
    # Load the Excel file from memory
    xls = pd.ExcelFile(BytesIO(file_bytes))


    # Function to process a single sheet
    def process_sheet_for_info(df):
        study_info = {
            "study_title": None,
            "study_id": None,
            "study_started": None,
            "study_ended": None,
            "study_status": False,
            "study_respondents": None,
            "study_keywords": []
        }

        for i in range(len(df)):
            row = df.iloc[i]
            if "Study Title" in str(row[0]):
                study_info["study_title"] = str(row[1]).strip()
            elif "Identification Number of the study" in str(row[0]):
                study_info["study_id"] = str(row[1]).strip()
            elif "Date when the study was run" in str(row[0]):
                date_range = str(row[1]).strip().replace("(", "").replace(")", "").split("-")
                study_info["study_started"] = date_range[0].strip()
                study_info["study_ended"] = date_range[1].strip() if len(date_range) > 1 else None
                study_info["study_status"] = date_range[1].strip() if len(date_range) > 1 else None
            elif "Number of respondents" in str(row[0]):
                study_info["study_respondents"] = int(row[1]) if pd.notna(row[1]) else None
            elif "Keywords" in str(row[0]):
                study_info["study_keywords"] = [kw.strip() for kw in str(row[1]).split(",") if kw.strip()]
        return study_info

    def process_sheet(df):
        try:
            mindset_cols = [col for col in df.columns if "Mindset" in str(col)]
            gender_cols = [col for col in df.columns if "Male" in str(col) or "Female" in str(col)]
            age_cols = [col for col in df.columns if any(age in str(col) for age in ["13 - 17", "18 - 24", "25 - 34", "35 - 44", "45 - 54", "55 - 64", "65+"])]
            total_col = "Total"

            prelim_answer_cols = [col for col in df.columns[3:] if col not in mindset_cols + gender_cols + age_cols + [total_col] and "Unnamed" not in str(col)]
            
            print("in parsesheet base  size  ")
            base_size = df.iloc[0, 2]
            print("completed parsesheet base  size  ")

            output_json = {
                "Base Size": base_size,
                "Questions": []
            }

            current_question = None
            for index, row in df.iterrows():
                if isinstance(row[1], str) and row[1].startswith("Question"):
                    current_question = {
                        "Question": row[1].split(": ", 1)[1] if ": " in row[1] else row[1],
                        "options": []
                    }
                    output_json["Questions"].append(current_question)
                elif isinstance(row[0], str) and current_question:
                    total=row[total_col] if total_col in df.columns else None
                    if str(total)=="nan":
                        total=None
                    print(total)
                    option_data = {
                        "optiontext": row[1],
                        "Total": total,
                        "Mindsets": [],
                        "Gender Segments": {},
                        "Age Segments": {},
                        "Prelim-Answer Segments": []
                    }

                    for col in mindset_cols:
                        option_data["Mindsets"].append({col: row[col] if str(row[col])!="nan" else None}) 
                    for col in gender_cols:
                        option_data["Gender Segments"][col] = row[col] if str(row[col])!="nan" else None
                    for col in age_cols:
                        option_data["Age Segments"][col] = row[col] if str(row[col])!="nan" else None
                    for col in prelim_answer_cols:
                        
                        option_data["Prelim-Answer Segments"].append({col: row[col] if str(row[col])!="nan" else None}) 

                    current_question["options"].append(option_data)

            return output_json
        except:
            return {}

    print(xls.sheet_names)
    study_info = process_sheet_for_info(xls.parse(sheet_name="Information Block", header=None).dropna(how="all"))
    print("study_info :- donee  ")
    print("study_info :- donee  ")
    print("study_info :- donee  ")
    print("study_info :- donee  ")
    # output_data = {sheet: process_sheet(xls.parse(sheet)) for sheet in xls.sheet_names[5:]}
    # Extract base values for all sheets
    def extract_base_values(df):
        base_values = {}
        for col in df.columns[3:]:
            base_values[col] = df.iloc[0][col] if not pd.isna(df.iloc[0][col]) else None
        return base_values
    
    
    output_data = {}
    for sheet in xls.sheet_names[5:]:
        try:
            df = xls.parse(sheet)
            output_data[sheet] = {
                "Base Values": extract_base_values(df),
                "Data": process_sheet(df)
            }
            
            
        except Exception as error :
            # print(error)
            output_data[sheet]={
                "error": error
            }
            
        # print(output_data)



    flat_output_data=flatten_all_segments(output_data)
    ffff=open("flat_output.json","w")
    ffff.write(json.dumps(flat_output_data, indent=4))
    ffff.close()
    return {"_id":str(uuid.uuid4()),
            "studyTitle":study_info['study_title'],
            "hasPpt":False,
            "studyStarted":study_info['study_started'],
            "studyEnded":study_info['study_ended'],
            "studyStatus":study_info['study_status'],
            "studyRespondents":study_info['study_respondents'],
            "studyKeywords":study_info['study_keywords'],
            "studyData":output_data,
            "studyData_flat":flat_output_data,
            "studySummarizationData":extract_summarizer_data(file_bytes),
            "mindsetsSegmentationsData": generate_segment_percentages(file_bytes)
            }


def generate_final_insights(json_data):
    insights = {}
    data_type = next(iter(json_data))
    
    # Handle Mindset data differently
    if 'Mindsets' in data_type:
        return generate_separated_mindset_insights(json_data)
    
    questions = json_data[data_type]['Data']['Questions']
    for i, q in enumerate(questions, 1):
        q_key = f"q{i}"
        insights[q_key] = []
        
        if 'Overall' in data_type:
            opts = sorted(q['options'], key=lambda x: x['Total'], reverse=True)
            insights[q_key].extend([
                f"Preferred highest: {opts[0]['Total']} ({opts[0]['optiontext']})",
                f"Preferred lowest: {opts[-1]['Total']} (Least preferred)",
                f"Preferred average: {sum(opt['Total'] for opt in opts)/len(opts):.1f} (Average)",
                f"Summary: {opts[0]['Total']} performed best in {opts[0]['optiontext']} (Excellent preferred)"
            ])
        
        elif 'Age' in data_type:
            age_stats = {}
            for opt in q['options']:
                for age, score in opt['Age Segments'].items():
                    if age not in age_stats or score > age_stats[age]:
                        age_stats[age] = score
            top_age, top_score = max(age_stats.items(), key=lambda x: x[1])
            insights[q_key].extend([
                f"Preferred highest: {top_score} (Age {top_age})",
                f"Preferred lowest: {min(age_stats.values())}",
                f"Preferred average: {sum(age_stats.values())/len(age_stats):.1f} (All ages)",
                f"Summary: Age {top_age} with score {top_score} performed best (Excellent preferred)"
            ])
        
        elif 'Gender' in data_type:
            gender_stats = {}
            for opt in q['options']:
                for gender, score in opt['Gender Segments'].items():
                    if gender not in gender_stats or score > gender_stats[gender]:
                        gender_stats[gender] = score
            top_gender, top_score = max(gender_stats.items(), key=lambda x: x[1])
            insights[q_key].extend([
                f"Preferred highest: {top_score} ({top_gender})",
                f"Preferred lowest: {min(gender_stats.values())}",
                f"Preferred average: {sum(gender_stats.values())/len(gender_stats):.1f} (Both genders)",
                f"Summary: {top_gender} with score {top_score} performed best (Excellent preferred)"
            ])
        
        elif 'Prelim' in data_type:
            segment_stats = {}
            for seg in [k for k in json_data[data_type]['Base Values'] if not str(k).startswith('Unnamed')]:
                max_score = max(
                    (s[seg] for opt in q['options'] 
                     for s in opt.get('Prelim-Answer Segments', []) 
                     if seg in s), 
                    default=0
                )
                if max_score > 0:
                    segment_stats[seg] = max_score
            top_segment, top_score = max(segment_stats.items(), key=lambda x: x[1])
            insights[q_key].extend([
                f"Preferred highest: {top_score} ('{top_segment}')",
                f"Preferred lowest: {min(segment_stats.values())}",
                f"Preferred average: {sum(segment_stats.values())/len(segment_stats):.1f} (All segments)",
                f"Summary: '{top_segment}' with score {top_score} performed best (Excellent preferred)"
            ])
    
    return insights

def generate_separated_mindset_insights(json_data):
    insights = {
        "2_mindset": {},
        "3_mindset": {}
    }
    
    # Find the mindset key (could be B, T, or R)
    mindset_key = next((k for k in json_data.keys() if 'Mindsets' in k), None)
    if not mindset_key:
        return insights  # Return empty structure if no mindset data found
    
    questions = json_data[mindset_key]['Data']['Questions']
    base_values = json_data[mindset_key]['Base Values']
    
    # Separate mindset groups
    mindset_2 = [k for k in base_values if 'Mindset' in k and 'of 2' in k]  # 1 of 2, 2 of 2
    mindset_3 = [k for k in base_values if 'Mindset' in k and 'of 3' in k]  # 1 of 3, 2 of 3, 3 of 3
    
    for i, question in enumerate(questions, 1):
        q_key = f"q{i}"
        question_text = question.get("Question", "")
        options = question.get("options", [])
        
        # Initialize question entries for both groups
        insights["2_mindset"][q_key] = []
        insights["3_mindset"][q_key] = []
        
        # Analyze 2-mindset group
        if mindset_2:
            group_scores = {}
            for m in mindset_2:
                max_score = max(
                    (s[m] for opt in options 
                     for s in opt.get('Mindsets', []) 
                     if m in s), 
                    default=0
                )
                if max_score > 0:
                    group_scores[m] = max_score
            
            if group_scores:
                top_mindset, top_score = max(group_scores.items(), key=lambda x: x[1])
                insights["2_mindset"][q_key].extend([
                    f"Preferred highest: {top_score} ({top_mindset})",
                    f"Preferred lowest: {min(group_scores.values())}",
                    f"Preferred average: {sum(group_scores.values())/len(group_scores):.1f}",
                    f"Summary: {top_mindset} performed best in this group"
                ])
        
        # Analyze 3-mindset group
        if mindset_3:
            group_scores = {}
            for m in mindset_3:
                max_score = max(
                    (s[m] for opt in options 
                     for s in opt.get('Mindsets', []) 
                     if m in s), 
                    default=0
                )
                if max_score > 0:
                    group_scores[m] = max_score
            
            if group_scores:
                top_mindset, top_score = max(group_scores.items(), key=lambda x: x[1])
                insights["3_mindset"][q_key].extend([
                    f"Preferred highest: {top_score} ({top_mindset})",
                    f"Preferred lowest: {min(group_scores.values())}",
                    f"Preferred average: {sum(group_scores.values())/len(group_scores):.1f}",
                    f"Summary: {top_mindset} performed best in this group"
                ])
    
    return insights

# def flatten_all_segments(raw_data):
#     structured_output = {}

#     for segment_key, segment_data in raw_data.items():
#         print(segment_key)
#         base_values = segment_data.get("Base Values", {})
#         questions = segment_data.get("Data", {}).get("Questions", [])

#         if not questions:
#             continue

#         # Check if this is an "Overall" segment (no inner segments)
#         is_overall_segment = (
#             "Overall" in segment_key or
#             (not base_values and all(
#                 not opt.get("Age Segments") and
#                 not opt.get("Gender Segments") and
#                 not opt.get("Mindsets") and
#                 not opt.get("Prelim-Answer Segments")
#                 for q in questions for opt in q.get("options", [])
#             ))
#         )

#         flattened = []

#         if is_overall_segment:
#             # Simple flattening with only question, option, total
#             for question in questions:
#                 q_text = question.get("Question", "")
#                 for option in question.get("options", []):
#                     row = {
#                         "Question": q_text,
#                         "Option": option.get("optiontext", ""),
#                         "Overall": option.get("Total", 0)
#                     }
#                     flattened.append(row)

#         else:
#             # Clean valid base segments for column naming
#             valid_segments = [k for k in base_values if k and not str(k).startswith("Unnamed")]

#             for question in questions:
#                 q_text = question.get("Question", "")
#                 for option in question.get("options", []):
#                     row = {
#                         "Question": q_text,
#                         "Option": option.get("optiontext", ""),
#                         "Overall": option.get("Total", 0)
#                     }

#                     # Loop through all potential segment types
#                     for segment_type_key, segment_value in option.items():
#                         if isinstance(segment_value, dict):
#                             for seg, val in segment_value.items():
#                                 col = f"{seg} ({int(base_values.get(seg, 0))})"
#                                 row[col] = val
#                         elif isinstance(segment_value, list):
#                             for item in segment_value:
#                                 if isinstance(item, dict):
#                                     for seg, val in item.items():
#                                         col = f"{seg} ({int(base_values.get(seg, 0))})"
#                                         row[col] = val

#                     # Fill in missing segments with 0
#                     for seg in valid_segments:
#                         col = f"{seg} ({int(base_values[seg])})"
#                         if col not in row:
#                             row[col] = 0

#                     flattened.append(row)

#         structured_output[segment_key] = flattened

#     return structured_output

def flatten_all_segments(raw_data):
    structured_output = {}

    for segment_key, segment_data in raw_data.items():
        base_values = segment_data.get("Base Values", {})
        questions = segment_data.get("Data", {}).get("Questions", [])

        if not questions:
            continue

        # Check if this is an "Overall" segment
        is_overall_segment = (
            "Overall" in segment_key or
            (not base_values and all(
                not opt.get("Age Segments") and
                not opt.get("Gender Segments") and
                not opt.get("Mindsets") and
                not opt.get("Prelim-Answer Segments")
                for q in questions for opt in q.get("options", [])
            ))
        )
        # Check if this is a Mindsets segment that needs splitting
        is_mindsets_segment = "Mindsets" in segment_key and not is_overall_segment

        flattened = []
        flattened_2_mindsets = []
        flattened_3_mindsets = []

        if is_overall_segment:
            # Simple flattening with only question, option, total
            for question in questions:
                q_text = question.get("Question", "")
                for option in question.get("options", []):
                    row = {
                        "Question": q_text,
                        "Option": option.get("optiontext", ""),
                        "Overall": option.get("Total", 0)
                    }
                    flattened.append(row)

        else:
            # Clean valid base segments for column naming
            valid_segments = [k for k in base_values if k and not str(k).startswith("Unnamed")]

            for question in questions:
                q_text = question.get("Question", "")
                for option in question.get("options", []):
                    row = {
                        "Question": q_text,
                        "Option": option.get("optiontext", ""),
                        "Overall": option.get("Total", 0)
                    }
                    
                    row_2_mindsets = {
                        "Question": q_text,
                        "Option": option.get("optiontext", ""),
                        "Overall": option.get("Total", 0)
                    }
                    
                    row_3_mindsets = {
                        "Question": q_text,
                        "Option": option.get("optiontext", ""),
                        "Overall": option.get("Total", 0)
                    }

                    # Loop through all potential segment types
                    for segment_type_key, segment_value in option.items():
                        if isinstance(segment_value, dict):
                            for seg, val in segment_value.items():
                                col = f"{seg} ({int(base_values.get(seg, 0))})"
                                row[col] = val
                                # Split mindsets into 2 and 3 groups
                                if is_mindsets_segment:
                                    if "1 of 2" in seg or "2 of 2" in seg:
                                        row_2_mindsets[col] = val
                                    elif "1 of 3" in seg or "2 of 3" in seg or "3 of 3" in seg:
                                        row_3_mindsets[col] = val
                        elif isinstance(segment_value, list):
                            for item in segment_value:
                                if isinstance(item, dict):
                                    for seg, val in item.items():
                                        col = f"{seg} ({int(base_values.get(seg, 0))})"
                                        row[col] = val
                                        # Split mindsets into 2 and 3 groups
                                        if is_mindsets_segment:
                                            if "1 of 2" in seg or "2 of 2" in seg:
                                                row_2_mindsets[col] = val
                                            elif "1 of 3" in seg or "2 of 3" in seg or "3 of 3" in seg:
                                                row_3_mindsets[col] = val

                    # Fill in missing segments with 0
                    for seg in valid_segments:
                        col = f"{seg} ({int(base_values[seg])})"
                        if col not in row:
                            row[col] = 0
                        if is_mindsets_segment:
                            if "1 of 2" in seg or "2 of 2" in seg:
                                if col not in row_2_mindsets:
                                    row_2_mindsets[col] = 0
                            elif "1 of 3" in seg or "2 of 3" in seg or "3 of 3" in seg:
                                if col not in row_3_mindsets:
                                    row_3_mindsets[col] = 0

                    flattened.append(row)
                    if is_mindsets_segment:
                        flattened_2_mindsets.append(row_2_mindsets)
                        flattened_3_mindsets.append(row_3_mindsets)

        if is_mindsets_segment:
            # Add both versions to the output
            structured_output[f"{segment_key}_2"] = flattened_2_mindsets
            structured_output[f"{segment_key}_3"] = flattened_3_mindsets
            # Also keep the original combined version
            structured_output[segment_key] = flattened
        else:
            structured_output[segment_key] = flattened

    return structured_output


def get_ppt(study_id,token):
   
    # # Set Chrome options
    # chrome_options = Options()
    # chrome_options = Options()
    # chrome_options.add_argument("--headless")  # Run in headless mode
    # chrome_options.add_argument("--start-maximized")  # Start maximized
    # chrome_options.add_argument("--window-size=1920,1080")  # Set window size (optional)
    # # unique_user_data_dir = f"/tmp/chrome-profile-{uuid.uuid4()}"
    # # chrome_options.add_argument(f"--user-data-dir={unique_user_data_dir}")

    # # Initialize WebDriver
    # driver = webdriver.Chrome(options=chrome_options)

    def get_driver():
        chrome_options = Options()

        # Create a unique temp directory for each session
        temp_profile_dir = tempfile.mkdtemp(prefix="chrome-profile-")

        # Use it as Chrome's user data dir
        chrome_options.add_argument(f"--user-data-dir={temp_profile_dir}")

        # Optional: headless mode
        chrome_options.add_argument("--headless=new")  # new headless mode (Chrome 109+)
        chrome_options.add_argument("--no-sandbox")
        chrome_options.add_argument("--window-size=1920,1080")  # Set window size (optional)
        chrome_options.add_argument("--disable-dev-shm-usage")

        driver = webdriver.Chrome(options=chrome_options)

        return driver, temp_profile_dir
    
    driver, temp_profile = get_driver()


    url = f"https://studiesapi.tikuntech.com/mf2/study/{study_id}"
    headers = {
        "Authorization": f"Bearer {token}"
    }
    print(headers)
    response = requests.get(url, headers=headers)

    if response.status_code == 200:
        print(response.json())  # Print the response JSON if successful
    else:
        print(f"Error: {response.status_code}, {response.text}")


    study_row_data=response.json()

    # Convert Python dictionary to JSON string
    USER_JSON = study_row_data['study']['studyCreatedBy']['user']



    # Step 1: Open the login page main
    # driver.get("https://mindgenome.org")



  # Step 1: Open the login page test server 
    driver.get("https://tikunstudies-test.netlify.app")
    time.sleep(5)  # Wait for page to fully load
    USER_JSON_STRING = json.dumps(USER_JSON)
    # Step 2: Inject Bearer Token and User Data into localStorage
    driver.execute_script(f"localStorage.setItem('authToken', '{token}');")
    driver.execute_script(f"localStorage.setItem('user', '{USER_JSON_STRING}');")

    print("✅ Token and user data injected into localStorage.")

    # Step 3: Confirm values are stored
    stored_token = driver.execute_script("return localStorage.getItem('authToken');")
    stored_user = driver.execute_script("return localStorage.getItem('user');")
    print("🔍 Stored Token:", stored_token)
    print("🔍 Stored User Data:", stored_user)

    # main server 
    # driver.get(f"https://mindgenome.org/study/{study_id}")


    # test server 
    driver.get(f"https://tikunstudies-test.netlify.app/study/{study_id}")
    time.sleep(10)  # Allow the page to load

    # Keep browser open for testing
    # input("Press Enter to close...")

    actions = ActionChains(driver)
    ppt = Presentation()
    study_title = driver.find_element(
        By.XPATH, '//*[@id="page-content"]/div/div[1]/h1'
    ).text


    # Modern Title Slide Design
    ppt_title_slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Blank layout

    # Clean background with subtle texture
    background = ppt_title_slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0x2A, 0x56, 0x8E)  # Dark Blue

    # Add a modern header bar
    header_bar = ppt_title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, ppt.slide_width, Inches(0.8)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x6E)  # Darker Blue

    # Main title with modern typography
    title_box = ppt_title_slide.shapes.add_textbox(
        Inches(1), Inches(2), 
        ppt.slide_width - Inches(2), Inches(2))
    title = title_box.text_frame
    title.paragraphs[0].text = study_title
    title.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
    title.paragraphs[0].font.size = Pt(42)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.name = 'Segoe UI'
    title.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Subtitle with accent color
    subtitle_box = ppt_title_slide.shapes.add_textbox(
        Inches(1), Inches(3.5),
        ppt.slide_width - Inches(2), Inches(1)
    )
    subtitle = subtitle_box.text_frame
    subtitle.paragraphs[0].text = "Automated Insights Report"
    subtitle.paragraphs[0].font.color.rgb = RGBColor(0xFD, 0x7E, 0x14)  # Orange
    subtitle.paragraphs[0].font.size = Pt(24)
    subtitle.paragraphs[0].font.name = 'Segoe UI Light'
    subtitle.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Add date and version info
    date_box = ppt_title_slide.shapes.add_textbox(
        Inches(1), ppt.slide_height - Inches(1),
        Inches(4), Inches(0.5)
    )
    date = date_box.text_frame
    date.paragraphs[0].text = datetime.datetime.now().strftime("%B %Y")
    date.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)  # Light Gray
    date.paragraphs[0].font.size = Pt(12)

    # Professional Overview Slide
    overview_slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Blank layout

    # Light background with subtle pattern
    overview_slide.background.fill.solid()
    overview_slide.background.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)

    # Title with modern styling
    title_box = overview_slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(6), Inches(1)
    )
    title = title_box.text_frame
    title.paragraphs[0].text = "STUDY OVERVIEW"
    title.paragraphs[0].font.color.rgb = RGBColor(0x2A, 0x56, 0x8E)  # Dark Blue
    title.paragraphs[0].font.size = Pt(28)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.name = 'Segoe UI'

    # Content box with clean shadow
    content_box = overview_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.5),
        Inches(12), Inches(4.5)
    )
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
    content_box.line.color.rgb = RGBColor(220, 220, 220)
    content_box.shadow.inherit = True

    # Professional content layout
    content = overview_slide.shapes.add_textbox(
        Inches(1), Inches(1.8),
        Inches(11), Inches(4)
    )
    tf = content.text_frame
    tf.word_wrap = True

    # Status item
    p = tf.add_paragraph()
    p.text = "• Status"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2A, 0x56, 0x8E)
    p = tf.add_paragraph()
    p.text = driver.find_element(By.XPATH,'//*[@id="page-content"]/div/div[2]/p[1]/span').text
    p.font.size = Pt(14)
    p.space_after = Inches(0.2)

    # Surveys Started item
    p = tf.add_paragraph()
    p.text = "• Surveys Started"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2A, 0x56, 0x8E)
    p = tf.add_paragraph()
    p.text = driver.find_element(By.XPATH,'//*[@id="page-content"]/div/div[2]/p[2]/span').text
    p.font.size = Pt(14)
    p.space_after = Inches(0.2)

    # Surveys Completed item
    p = tf.add_paragraph()
    p.text = "• Surveys Completed"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2A, 0x56, 0x8E)
    p = tf.add_paragraph()
    p.text = driver.find_element(By.XPATH,'//*[@id="page-content"]/div/div[2]/p[3]').text.replace("Surveys Completed:","")
    p.font.size = Pt(14)

    # Add a subtle footer line
    footer_line = overview_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), ppt.slide_height - Inches(0.8),
        Inches(12), Inches(0.02)
    )
    footer_line.fill.solid()
    footer_line.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    output_dir = "screenshots"
    os.makedirs(output_dir, exist_ok=True)

    categories = ["Overall", "Age", "Gender", "Prelim", "2 Market Segments", "3 Market Segments"]
    views = ["Top Down", "Bottom Up", "Response Time"]

    visuals={
        # "table":'//*[@id="page-content"]/div/div[3]/div[1]/button[1]',
        "heatmap":'//*[@id="page-content"]/div/div[3]/div[1]/button[2]',
        "charts":'//*[@id="page-content"]/div/div[3]/div[1]/button[3]'
    }


    from selenium.common.exceptions import NoSuchElementException

    useful_xpaths={
        "Overall":'//*[@id="radix-:r0:-content-overall"]/div/div',
        "Age":'//*[@id="radix-:r0:-content-age"]/div/div',
        "Gender":'//*[@id="radix-:r0:-content-gender"]/div/div',
        "Prelim":'//*[@id="radix-:r0:-content-prelim"]/div/div',
        "2 Market Segments":'//*[@id="radix-:r0:-content-2 market segments"]/div/div',
        "3 Market Segments":'//*[@id="radix-:r0:-content-3 market segments"]/div/div'
    }


    def crop_and_resize_image(image_path, output_path, crop_top=0, crop_sides=0):
        """Crop and resize image while maintaining aspect ratio"""
        try:
            with Image.open(image_path) as img:
                width, height = img.size
                # Crop from top and sides
                left = crop_sides
                top = crop_top
                right = width - crop_sides
                bottom = height
                img = img.crop((left, top, right, bottom))
                img.save(output_path)
            return True
        except Exception as e:
            print(f"Error processing {image_path}: {str(e)}")
            return False




    def create_presentation(ppt, questions, categories, views, base_path):
        # Color scheme
        PRIMARY_COLOR = RGBColor(42, 86, 142)    # Dark Blue
        SECONDARY_COLOR = RGBColor(253, 126, 20)  # Orange
        ACCENT_COLOR = RGBColor(40, 167, 69)      # Green
        BACKGROUND_COLOR = RGBColor(248, 249, 250) # Light Gray

        # Set master slide background
        ppt.slide_master.background.fill.solid()
        ppt.slide_master.background.fill.fore_color.rgb = BACKGROUND_COLOR

        # Slide dimensions (16:9 aspect ratio)
        ppt.slide_width = Inches(13.33)
        ppt.slide_height = Inches(7.5)

        blank_slide_layout = ppt.slide_layouts[6]

        # Helper function to add bordered images
        def add_bordered_image(slide, left, top, width, height):
            border = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left - Inches(0.05), top - Inches(0.05),
                width + Inches(0.1), height + Inches(0.1)
            )
            border.fill.solid()
            border.fill.fore_color.rgb = RGBColor(255, 255, 255)
            border.line.color.rgb = PRIMARY_COLOR
            border.line.width = Pt(1.5)
            return border

        # Helper function to safely add text with auto-fitting
        def add_safe_text(text_frame, text, max_font=18, min_font=10, bold=False, color=None):
            # First try with max font size
            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = text if len(text) <= 150 else text[:140] + "..."
            run.font.size = Pt(max_font)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color
            
            return p

        for category in categories:
            for view in views:
                if category=="Overall" and view=="Top Down":
                    insights=generate_final_insights({"(T) Overall":study_row_data["study"]["studyData"]["(T) Overall"]})
                elif category=="Overall" and view=="Bottom Up":
                    insights=generate_final_insights({"(B) Overall":study_row_data["study"]["studyData"]["(B) Overall"]})
                elif category=="Overall" and view=="Response Time":
                    insights=generate_final_insights({"(R) Overall":study_row_data["study"]["studyData"]["(R) Overall"]})
                
                elif category=="Age" and view=="Top Down":
                    insights=generate_final_insights({"(T) Age segments":study_row_data["study"]["studyData"]["(T) Age segments"]})
                
                elif category=="Age" and view=="Bottom Up":
                    insights=generate_final_insights({"(B) Age segments":study_row_data["study"]["studyData"]["(B) Age segments"]})
                elif category=="Age" and view=="Response Time":
                
                    insights=generate_final_insights({"(R) Age segments":study_row_data["study"]["studyData"]["(R) Age segments"]})
                if category=="Gender" and view=="Top Down":
                    
                    insights=generate_final_insights({"(T) Gender segments":study_row_data["study"]["studyData"]["(T) Gender segments"]})
                elif category=="Gender" and view=="Bottom Up":
                    insights=generate_final_insights({"(B) Gender segments":study_row_data["study"]["studyData"]["(B) Gender segments"]})
                elif category=="Gender" and view=="Response Time":
                    insights=generate_final_insights({"(R) Gender segments":study_row_data["study"]["studyData"]["(R) Gender segments"]})
                
                elif category=="Prelim" and view=="Top Down":
                    insights=generate_final_insights({"(T) Prelim-answer segments":study_row_data["study"]["studyData"]["(T) Prelim-answer segments"]})
                elif category=="Prelim" and view=="Bottom Up":
                    insights=generate_final_insights({"(B) Prelim-answer segments":study_row_data["study"]["studyData"]["(B) Prelim-answer segments"]})
                elif category=="Prelim" and view=="Response Time":
                    insights=generate_final_insights({"(R) Prelim-answer segments":study_row_data["study"]["studyData"]["(R) Prelim-answer segments"]})
                
                elif category=="2 Market Segments" and view=="Top Down":
                    insights=generate_final_insights({"(T) Mindsets":study_row_data["study"]["studyData"]["(T) Mindsets"]})["2_mindset"]
                elif category=="2 Market Segments" and view=="Bottom Up":
                    insights=generate_final_insights({"(B) Mindsets":study_row_data["study"]["studyData"]["(B) Mindsets"]})["2_mindset"]
                elif category=="2 Market Segments" and view=="Response Time":
                    insights=generate_final_insights({"(R) Mindsets":study_row_data["study"]["studyData"]["(R) Mindsets"]})["2_mindset"]
                    
                    
                elif category=="3 Market Segments" and view=="Top Down":
                    insights=generate_final_insights({"(T) Mindsets":study_row_data["study"]["studyData"]["(T) Mindsets"]})["3_mindset"]
                elif category=="3 Market Segments" and view=="Bottom Up":
                    insights=generate_final_insights({"(B) Mindsets":study_row_data["study"]["studyData"]["(B) Mindsets"]})["3_mindset"]
                elif category=="3 Market Segments" and view=="Response Time":
                    insights=generate_final_insights({"(R) Mindsets":study_row_data["study"]["studyData"]["(R) Mindsets"]})["3_mindset"]
                
                
                
                print(insights)
                # Title slide
                title_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
                title = title_slide.shapes.title
                title.text = f"{category} ({view})"
                title.text_frame.paragraphs[0].font.size = Pt(36)
                title.text_frame.paragraphs[0].font.bold = True
                title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
                title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                subtitle = title_slide.placeholders[1]
                subtitle.text = "Consumer Insights Analysis"
                subtitle.text_frame.paragraphs[0].font.size = Pt(24)
                subtitle.text_frame.paragraphs[0].font.color.rgb = SECONDARY_COLOR
                subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                # Decorative line
                line = title_slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    0, Inches(2.5),
                    ppt.slide_width, Inches(0.1)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = ACCENT_COLOR
                line.line.color.rgb = ACCENT_COLOR

                for q_idx in range(1, 5):
                    slide = ppt.slides.add_slide(blank_slide_layout)
                    
                    # Header bar
                    header = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        0, 0,
                        ppt.slide_width, Inches(0.8)
                    )
                    header.fill.solid()
                    header.fill.fore_color.rgb = PRIMARY_COLOR

                    # Header text (with safe handling)
                    header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(4), Inches(0.6))
                    tf = header_text.text_frame
                    add_safe_text(tf, f"{category} | {view}", 
                                max_font=16, min_font=12,
                                bold=True, 
                                color=RGBColor(255, 255, 255))

                    # Content positioning
                    content_top = Inches(1.0)
                    left_col = Inches(0.5)
                    col_width = Inches(6)
                    right_col = left_col + col_width + Inches(0.2)
                    right_width = Inches(5.3)
                    max_content_height = ppt.slide_height - content_top - Inches(1.0)

                    # Standard image dimensions
                    chart_height = Inches(3)
                    heatmap_height = Inches(3)
                    image_gap = Inches(0.3)

                    # Chart image with border
                    chart_path = os.path.join(base_path, f"{category}_charts_Q{q_idx}_{view}.png")
                    if os.path.exists(chart_path):
                        add_bordered_image(slide, left_col, content_top, col_width, chart_height)
                        slide.shapes.add_picture(
                            chart_path, left_col, content_top,
                            width=col_width, height=chart_height
                        )

                    # Heatmap image with border
                    heatmap_path = os.path.join(base_path, f"{category}_heatmap_Q{q_idx}_{view}.png")
                    heatmap_top = content_top + chart_height + image_gap
                    if os.path.exists(heatmap_path):
                        add_bordered_image(slide, left_col, heatmap_top, col_width, heatmap_height)
                        slide.shapes.add_picture(
                            heatmap_path, left_col, heatmap_top,
                            width=col_width, height=heatmap_height
                        )

                    # Text content with borders
                    try:
                        # Content box dimensions
                        content_width = Inches(5.5)
                        content_left = right_col
                        content_height = min(Inches(4.8), max_content_height)
                        available_height = ppt.slide_height - content_top - Inches(0.5)
                        content_top = content_top + (available_height - content_height) / 2

                        # Create content box
                        content_box = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            content_left, content_top,
                            content_width, content_height
                        )
                        content_box.fill.solid()
                        content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                        content_box.line.color.rgb = ACCENT_COLOR
                        content_box.line.width = Pt(1.5)

                        # Text frame with padding
                        txBox = slide.shapes.add_textbox(
                            content_left + Inches(0.3),
                            content_top + Inches(0.2),
                            content_width - Inches(0.6),
                            content_height - Inches(0.1)
                        )
                        tf = txBox.text_frame
                        tf.word_wrap = True
                        
                        # Add question (simplified without add_run)
                        p = tf.add_paragraph()
                        p.text = questions[f"Q{q_idx}"][:100] + "..." if len(questions[f"Q{q_idx}"]) > 150 else questions[f"Q{q_idx}"]
                        p.alignment = PP_ALIGN.LEFT
                        p.font.size = Pt(16)
                        p.font.bold = True
                        p.font.color.rgb = PRIMARY_COLOR
                        p.space_after = Inches(0.1)

                        # Key Findings title
                        p = tf.add_paragraph()
                        p.text = "Key Findings:"
                        p.alignment = PP_ALIGN.LEFT
                        p.font.size = Pt(12)
                        p.font.bold = True
                        p.font.color.rgb = SECONDARY_COLOR
                        p.space_after = Inches(0.15)

                        # Findings as bullet points
                        for finding in insights[f"q{q_idx}"]:
                            p = tf.add_paragraph()
                            p.text = f"• {finding}"
                            p.alignment = PP_ALIGN.LEFT
                            p.space_before = Inches(0.15)
                            p.font.size = Pt(14)
                    
                        p = tf.add_paragraph()
                        p.text = "© mindgenome.org"
                        p.alignment = PP_ALIGN.CENTER
                        p.space_before = Inches(0.2)
                        p.font.size = Pt(12)
                        p.font.color.rgb = RGBColor(150, 150, 150)

                    except Exception as e:
                        print(f"Error creating content: {str(e)}")


        # Clean up temp files
        temp_files = [
            os.path.join(base_path, "temp_chart.png"),
            os.path.join(base_path, "temp_heatmap.png")
        ]
        for temp_file in temp_files:
            if os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except Exception as e:
                    print(f"Error removing temp file {temp_file}: {str(e)}")
        
        try:
            ppt.save(f"{study_title}.pptx")
            print(f"Successfully created presentation: {study_title}.pptx")
        except Exception as e:
            print(f"Error saving presentation: {str(e)}")

    def set_title_style(title_shape):
        title_shape.text_frame.text = title_shape.text_frame.text.upper()
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Loop through each category and check if the element exists
    questions={"Q1":"","Q2":"","Q3":"","Q4":""}
    for category in categories:
        try:
            element_category = driver.find_element(By.XPATH, f"//*[contains(text(), '{category}')]").click()
            
            print(f"✅ Found category: {category}")
            for visual_name, visual_xpath in visuals.items():
                try:
                    element_virual = driver.find_element(By.XPATH, visual_xpath).click()
                    print(f"✅ Found visual: {visual_name}")
                    all_divs=driver.find_elements(By.XPATH,useful_xpaths[category])
                    for idx, div in enumerate(all_divs,start=1):
                        for view in views:
                            try:
                                # Click on the tab
                                element__ = driver.find_element(By.XPATH, f"//*[contains(text(), '{view}')]").click()
                                driver.execute_script("document.body.style.cursor = 'default';")
                        
                                element = driver.find_element(By.XPATH, f'{useful_xpaths[category]}[{idx}]')
                                # time.sleep(0.15)
                                # print(f'{useful_xpaths[category]}[{idx}]/div[1]')
                                driver.execute_script("arguments[0].scrollIntoView();", element)  # Scroll to the tab
                                
                                # actions.move_to_element_with_offset(element, -100, -100).perform()  # Move mouse away
                                driver.execute_script("document.body.style.cursor = 'default';")  # Reset cursor
                                element_title=element.find_element(By.TAG_NAME,"h2")
                                element_title.click()
                                time.sleep(0.15)
                                screenshot_path = os.path.join(output_dir, f"{category}_{visual_name}_Q{idx}_{view}.png")
                                questions[f"Q{idx}"]=element_title.text
                                element.screenshot(screenshot_path)
                                print(f"✅ Captured {view} for {category}")
                            

                            except NoSuchElementException:
                                print(f"❌ Tab or element not found: {view}")
                            except :
                                print(f"⚠️ Click intercepted for tab: {view}")
                                        
                        
                        
                        
                        
                        
                        
                    
                except NoSuchElementException:
                    print(f"❌ Visual not found: {visual_name}")
        except NoSuchElementException:
            print(f"❌ Category not found: {category}")
    
    create_presentation(ppt,questions,categories,views,"./screenshots")      
    ppt.save(f"{study_title}.pptx")
    # ppt.save(f"test.pptx")
    with open(f"{study_title}.pptx", "rb") as file:
        ppturl = upload_to_azure_with_filename(file,f"{study_title}.pptx")
        print("Uploaded to:", ppturl)
    os.remove(f"{study_title}.pptx")
    print("✅ PowerPoint Report Generated: Study_Analysis_Report.pptx")
    
    STUDIES_collection.find_one_and_update({"_id":study_id},  {"$set": {"hasPpt": True, "pptUrl":ppturl}})
    driver.quit()

    shutil.rmtree(temp_profile, ignore_errors=True)  # Clean up
    # input("Press Enter to exit...")


def generate_segment_percentages(file_bytes: bytes):
    # === Load Excel file from bytes ===
    workbook = pd.ExcelFile(BytesIO(file_bytes))
    raw_df = workbook.parse("RawData")
    info_df = workbook.parse("Information Block", header=None)

    # === STEP 1: Keep first row per Panelist ===
    deduped_df = raw_df.drop_duplicates(subset="Panelist", keep="first")

    # === STEP 2: Map Gender & AgeGroup ===
    gender_map = {1: "Male", 2: "Female"}
    age_map = {
        1: "13-17", 2: "18-24", 3: "25-34", 4: "35-44",
        5: "45-54", 6: "55-64", 7: "65+"
    }

    deduped_df["Gender"] = deduped_df["Gender"].map(gender_map)
    deduped_df["AgeGroup"] = deduped_df["AgeGroup"].map(age_map)

    # === STEP 3: Extract segments from Information Block ===
    segments = []
    q_index = 1
    i = 0

    while i < len(info_df):
        line = str(info_df.iloc[i, 0]).strip()
        if line.startswith("Preliminary question"):
            question = info_df.iloc[i, 1]
            answers_line = str(info_df.iloc[i + 1, 1])
            attr_map = {}

            for ans in answers_line.split():
                if "=" in ans:
                    code, label = ans.split("=")
                    attr_map[int(code)] = label.strip()

            segments.append({
                "segment": question,
                "column": f"Quest{q_index}",
                "attrMap": attr_map
            })
            q_index += 1
            i += 2
        else:
            i += 1

    # === STEP 4: Add fixed segments (Gender, AgeGroup) ===
    segments.insert(0, {
        "segment": "AgeGroup",
        "column": "AgeGroup",
        "attrMap": {v: v for v in age_map.values()}
    })
    segments.insert(0, {
        "segment": "Gender",
        "column": "Gender",
        "attrMap": {"Male": "Male", "Female": "Female"}
    })

    # === STEP 5: Calculate mindset percentages ===
    total = len(deduped_df)
    mindset_counts = deduped_df["THREE MS"].value_counts().to_dict()

    result = []

    for seg in segments:
        segment = seg["segment"]
        column = seg["column"]
        attr_map = seg["attrMap"]

        for code, label in attr_map.items():
            matching = deduped_df[deduped_df[column].astype(str) == str(code)]
            row = {
                "Segment": segment,
                "Attribute": label,
                "%TotalRespondents": f"{round(len(matching) / total * 100)}%"
            }

            for m in range(1, 4):
                mind_match = matching[matching["THREE MS"] == m]
                ms_total = mindset_counts.get(m, 1)  # avoid div by 0
                row[f"%Mindset{m}"] = f"{round(len(mind_match) / ms_total * 100)}%"

            result.append(row)

    return result





# # # Load file and read bytes
# with open("04032025.Clean_.xlsx", "rb") as f:
#     file_bytes = f.read()

# # Call the function
# results = get_file_data_for_study(file_bytes)
